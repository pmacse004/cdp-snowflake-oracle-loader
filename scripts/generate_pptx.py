"""
generate_pptx.py
Generates: IBM_Bob_Stand_Deliver.pptx
8 slides — National Grid DSM / CDP Snowflake-to-Oracle Loader use case.
Run: python scripts/generate_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Colour Palette ───────────────────────────────────────────────
IBM_BLUE        = RGBColor(0x00, 0x43, 0xCE)
IBM_BLUE_LIGHT  = RGBColor(0xD0, 0xE2, 0xFF)
WHITE           = RGBColor(0xFF, 0xFF, 0xFF)
DARK            = RGBColor(0x16, 0x16, 0x16)
MUTED           = RGBColor(0x52, 0x52, 0x52)
GREEN           = RGBColor(0x19, 0x80, 0x38)
GREEN_LIGHT     = RGBColor(0xDE, 0xF7, 0xE9)
RED             = RGBColor(0xDA, 0x1E, 0x28)
RED_LIGHT       = RGBColor(0xFF, 0xD7, 0xD9)
YELLOW_LIGHT    = RGBColor(0xFF, 0xF1, 0xC2)
YELLOW          = RGBColor(0xF1, 0xC2, 0x1B)
ACCENT          = RGBColor(0x69, 0x29, 0xD0)
SLIDE_BG        = RGBColor(0xF4, 0xF4, 0xF4)
SLIDE_BG_DARK   = RGBColor(0x26, 0x26, 0x26)
NG_GREEN        = RGBColor(0x00, 0x6A, 0x2E)   # National Grid green
NG_GREEN_LIGHT  = RGBColor(0xCC, 0xEB, 0xD7)


# ── Helpers ──────────────────────────────────────────────────────

def add_solid_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height,
             fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_label(slide, text, left, top, width, height,
              font_size=12, bold=False, color=DARK,
              align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_multiline(slide, lines, left, top, width, height,
                  font_size=11, default_color=DARK, align=PP_ALIGN.LEFT):
    """lines: list of str  OR  (text, bold, color_or_None)"""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, bold, col = item, False, default_color
        else:
            text = item[0]
            bold = item[1] if len(item) > 1 else False
            col  = item[2] if (len(item) > 2 and item[2]) else default_color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = col
    return txBox


def add_footer(slide, text="National Grid DSM  |  CDP Snowflake Loader  |  IBM Bob Level 3  |  Althaf P M  |  2025"):
    add_rect(slide, 0, 7.44, 13.33, 0.06, fill_color=IBM_BLUE)
    add_label(slide, text, 0.2, 7.46, 12.9, 0.25,
              font_size=7, color=WHITE, align=PP_ALIGN.CENTER)


def divider(slide, top, color=IBM_BLUE_LIGHT, width=12.9):
    add_rect(slide, 0.2, top, width, 0.02, fill_color=color)


def script_bar(slide, time_range, script_text):
    add_rect(slide, 0.25, 6.88, 12.85, 0.38,
             fill_color=RGBColor(0x39, 0x39, 0x3A),
             line_color=RGBColor(0x52, 0x52, 0x52), line_width=Pt(0.5))
    add_label(slide,
              f"SCRIPT [{time_range}]  \"{script_text}\"",
              0.4, 6.9, 12.6, 0.35,
              font_size=8, italic=True,
              color=RGBColor(0xC6, 0xC6, 0xC6))


# ── Create Presentation ──────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE  (0:00 – 0:12)
# ════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
add_solid_background(s1, SLIDE_BG_DARK)
add_rect(s1, 0, 0, 13.33, 0.18, fill_color=NG_GREEN)

# Left blue panel
add_rect(s1, 0, 0.18, 5.8, 7.32, fill_color=IBM_BLUE)

# IBM Bob wordmark
add_label(s1, "IBM", 0.35, 0.5, 3.0, 1.1,
          font_size=60, bold=True, color=WHITE)
add_label(s1, "Bob", 0.35, 1.5, 3.0, 1.0,
          font_size=60, bold=False, color=IBM_BLUE_LIGHT)
add_rect(s1, 0.35, 2.52, 5.1, 0.04, fill_color=WHITE)

add_label(s1, "Level 3  \u00b7  Stand & Deliver", 0.35, 2.68, 5.1, 0.45,
          font_size=14, color=WHITE)

add_label(s1, "National Grid\nDemand Side Management (DSM)", 0.35, 3.25, 5.1, 1.35,
          font_size=26, bold=True, color=WHITE)

add_label(s1, "CDP Snowflake \u2192 DSM Oracle Loader", 0.35, 4.7, 5.1, 0.5,
          font_size=13, color=IBM_BLUE_LIGHT)

add_label(s1, "Althaf P M  |  IBM Technical Sales  |  2025",
          0.35, 6.6, 5.1, 0.45,
          font_size=10, color=IBM_BLUE_LIGHT)

# Right side
add_label(s1, "Client", 6.2, 0.9, 6.8, 0.35,
          font_size=10, bold=True, color=NG_GREEN_LIGHT)
add_label(s1, "National Grid", 6.2, 1.22, 6.8, 0.5,
          font_size=20, bold=True, color=WHITE)
add_label(s1, "20+ years of energy efficiency programmes\nacross Massachusetts & New York State",
          6.2, 1.72, 6.8, 0.65, font_size=11, color=RGBColor(0xC6, 0xC6, 0xC6))

add_rect(s1, 6.2, 2.45, 6.8, 0.02, fill_color=NG_GREEN_LIGHT)

add_label(s1, "Application", 6.2, 2.58, 6.8, 0.3,
          font_size=10, bold=True, color=NG_GREEN_LIGHT)
add_label(s1, "Demand Side Management (DSM)", 6.2, 2.88, 6.8, 0.45,
          font_size=16, bold=True, color=WHITE)
add_label(s1,
          "Web-based J2EE application managing energy\n"
          "efficiency programmes: Large C&I, SBS, Energy Wise,\n"
          "Residential, AM. Handles vendor payments, invoicing,\n"
          "scheduling, construction tracking via SAP & CSS.",
          6.2, 3.35, 6.8, 1.3, font_size=10,
          color=RGBColor(0xC6, 0xC6, 0xC6))

add_label(s1, "AI Maturity Stage", 6.2, 4.75, 3.0, 0.35,
          font_size=10, bold=True, color=IBM_BLUE_LIGHT)
add_rect(s1, 6.2, 5.1, 3.0, 0.85, fill_color=ACCENT)
add_label(s1, "DELEGATE", 6.2, 5.1, 3.0, 0.85,
          font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_label(s1, "Bob owns design &\nimplementation.\nEngineer reviews & approves.",
          9.35, 5.1, 3.6, 0.85, font_size=9,
          color=RGBColor(0xC6, 0xC6, 0xC6))

# Script note
add_rect(s1, 6.2, 6.1, 6.8, 1.1,
         fill_color=RGBColor(0x39, 0x39, 0x3A),
         line_color=RGBColor(0x52, 0x52, 0x52), line_width=Pt(0.5))
add_label(s1, "SCRIPT  [0:00 - 0:12]", 6.35, 6.15, 6.6, 0.28,
          font_size=9, bold=True, color=IBM_BLUE_LIGHT)
add_label(s1,
          '"Hi, I\'m Althaf, a Technical Seller at IBM. Today I want to show you '
          'how IBM Bob helped National Grid build a future-ready data integration '
          'layer for their Demand Side Management application \u2014 and the exact '
          'numbers that prove it."',
          6.35, 6.43, 6.6, 0.75,
          font_size=8.5, italic=True, color=RGBColor(0xE0, 0xE0, 0xE0))

add_footer(s1)


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — BUSINESS CONTEXT & THE CHALLENGE  (0:12 – 0:30)
# ════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
add_solid_background(s2, SLIDE_BG)
add_rect(s2, 0, 0, 13.33, 0.18, fill_color=NG_GREEN)

add_label(s2, "The Business Context", 0.3, 0.28, 9.0, 0.52,
          font_size=24, bold=True, color=NG_GREEN)
add_label(s2, "National Grid DSM \u2014 20+ years of energy efficiency, now facing a major system migration",
          0.3, 0.82, 12.5, 0.35, font_size=12, color=MUTED)
divider(s2, 1.2, color=NG_GREEN)

# Left: DSM programmes
add_rect(s2, 0.25, 1.3, 5.9, 4.4,
         fill_color=WHITE, line_color=NG_GREEN, line_width=Pt(1))
add_rect(s2, 0.25, 1.3, 5.9, 0.42, fill_color=NG_GREEN)
add_label(s2, "What DSM Manages", 0.4, 1.33, 5.7, 0.38,
          font_size=12, bold=True, color=WHITE)

programmes = [
    ("Large Commercial & Industrial",
     "Energy efficiency retrofits for large C&I customers.\nTracking, workflow, vendor payments."),
    ("SBS (Small Business Solutions)",
     "Programme delivery for small businesses.\nScheduling, invoicing, construction tracking."),
    ("Energy Wise / Residential",
     "Home energy efficiency programmes.\nLighting, HVAC upgrades, rebate management."),
    ("AM (Account Management)",
     "Key account management & regulatory reporting\nto MA and NY state regulators."),
]
for i, (prog, desc) in enumerate(programmes):
    ly = 1.85 + i * 0.95
    add_rect(s2, 0.38, ly, 5.64, 0.82,
             fill_color=NG_GREEN_LIGHT, line_color=NG_GREEN, line_width=Pt(0.5))
    add_label(s2, prog, 0.52, ly + 0.04, 5.4, 0.3,
              font_size=10, bold=True, color=NG_GREEN)
    add_label(s2, desc, 0.52, ly + 0.34, 5.4, 0.45,
              font_size=9, color=DARK)

# Right: The migration challenge
add_rect(s2, 6.55, 1.3, 6.5, 4.4,
         fill_color=WHITE, line_color=IBM_BLUE, line_width=Pt(1))
add_rect(s2, 6.55, 1.3, 6.5, 0.42, fill_color=IBM_BLUE)
add_label(s2, "The Migration Challenge", 6.7, 1.33, 6.3, 0.38,
          font_size=12, bold=True, color=WHITE)

add_label(s2, "Current State", 6.7, 1.88, 6.2, 0.3,
          font_size=10, bold=True, color=IBM_BLUE)
add_rect(s2, 6.7, 2.2, 6.2, 0.75,
         fill_color=IBM_BLUE_LIGHT, line_color=IBM_BLUE, line_width=Pt(0.5))
add_label(s2,
          "CSS Billing System  \u2192  ETL (DataStage/Informatica)  \u2192  CDI Database  \u2192  DSM",
          6.8, 2.27, 6.0, 0.55, font_size=9.5, color=DARK)

add_label(s2, "Future State (Planned)", 6.7, 3.1, 6.2, 0.3,
          font_size=10, bold=True, color=GREEN)
add_rect(s2, 6.7, 3.42, 6.2, 0.75,
         fill_color=GREEN_LIGHT, line_color=GREEN, line_width=Pt(0.5))
add_label(s2,
          "Kraken System  \u2192  CDP Snowflake  \u2192  NEW LOADER  \u2192  DSM Oracle Tables",
          6.8, 3.49, 6.0, 0.55, font_size=9.5, color=DARK)

add_label(s2, "The Prototype We Built with IBM Bob", 6.7, 4.3, 6.2, 0.35,
          font_size=10, bold=True, color=ACCENT)
add_rect(s2, 6.7, 4.65, 6.2, 0.82,
         fill_color=RGBColor(0xED, 0xE9, 0xFE), line_color=ACCENT, line_width=Pt(0.75))
add_label(s2,
          "Read customer data from CDP Snowflake and load into\n"
          "DSM Oracle tables \u2014 replacing the CDI dependency,\n"
          "ready for the CSS \u2192 Kraken migration.",
          6.8, 4.7, 6.0, 0.75, font_size=9.5, color=DARK)

# Bottom stat bar
add_rect(s2, 0.25, 5.82, 12.85, 0.62, fill_color=DARK)
add_label(s2,
          "DSM serves MA & NY regulated energy efficiency programmes  \u00b7  "
          "Connected to SAP, CSS, LDAP, Business Objects  \u00b7  "
          "Millions of dollars in vendor payments & customer rebates annually",
          0.4, 5.87, 12.6, 0.55,
          font_size=9, color=WHITE, align=PP_ALIGN.CENTER)

script_bar(s2, "0:12-0:30",
           "National Grid has run energy efficiency programmes for 20+ years across MA and NY. "
           "Their DSM application manages the full workflow \u2014 from vendor payments to regulatory reporting. "
           "The challenge: CSS, their billing system, is being replaced by Kraken. "
           "DSM needs a new data feed from CDP Snowflake. That is what we built with IBM Bob.")
add_footer(s2)


# ════════════════════════════════════════════════════════════════
# SLIDE 3 — WITHOUT IBM BOB (The Manual Approach)  (0:30 – 0:52)
# ════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
add_solid_background(s3, SLIDE_BG)
add_rect(s3, 0, 0, 13.33, 0.18, fill_color=RED)

add_label(s3, "Without IBM Bob \u2014 The Manual Approach", 0.3, 0.28, 10.0, 0.52,
          font_size=24, bold=True, color=RED)
add_label(s3,
          "How this Snowflake \u2192 DSM loader would have been built traditionally",
          0.3, 0.82, 12.0, 0.35, font_size=12, color=MUTED)
divider(s3, 1.2, color=RED)

pains = [
    ("Requirements &\nMapping Analysis",
     "3-4 days: interview DSM team, reverse-engineer CDI schema, "
     "manually document Snowflake-to-DSM column mappings in spreadsheets. "
     "No traceability to code."),
    ("Schema &\nDDL Design",
     "2-3 days: hand-write Oracle DDL for staging/target tables, "
     "Flyway migrations. Prone to missing indexes, CHECK constraints. "
     "Environment drift across dev/test/prod."),
    ("ETL Pipeline\nDevelopment",
     "10-14 days: build Spring Batch readers (Snowflake JDBC), "
     "processors (transformation logic per entity), writers (Oracle MERGE), "
     "watermark strategy, skip/retry policy \u2014 all from scratch."),
    ("Snowflake\nProvisioning",
     "3-4 days: manual SQL scripts with no pre-flight safety checks, "
     "no idempotent re-run guarantee. Errors on partial failure require "
     "manual cleanup."),
    ("REST API &\nDashboard",
     "5-7 days: REST endpoints for job control, status, errors, "
     "reconciliation + React dashboard. Separate frontend resource needed."),
    ("Testing &\nDocumentation",
     "4-5 days: unit tests deferred, coverage gaps, "
     "docs written after the fact and immediately stale. "
     "Onboarding a new dev: 2-3 weeks."),
]
col_xs = [0.25, 4.55, 8.85]
for i, (title, desc) in enumerate(pains):
    col = i % 3
    row = i // 3
    lx = col_xs[col]
    ly = 1.32 + row * 2.12
    add_rect(s3, lx, ly, 4.05, 1.95,
             fill_color=RED_LIGHT, line_color=RED, line_width=Pt(0.75))
    add_label(s3, title, lx + 0.12, ly + 0.08, 3.82, 0.45,
              font_size=10.5, bold=True, color=RED)
    add_label(s3, desc, lx + 0.12, ly + 0.54, 3.82, 1.3,
              font_size=9, color=DARK)

add_rect(s3, 0.25, 5.58, 12.85, 0.72, fill_color=RED)
add_label(s3,
          "TOTAL:  6-8 WEEKS   |   3 people (engineer + BA + QA)   |   "
          "High mapping drift   |   No audit trail   |   Tight CSS/CDI coupling retained",
          0.4, 5.64, 12.6, 0.6,
          font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

script_bar(s3, "0:30-0:52",
           "Without Bob: 3-4 days just mapping Snowflake columns to DSM tables. "
           "A week building Spring Batch alone. No traceability between specs and code. "
           "6-8 weeks, 3 people, and the spec-to-code drift is a real risk on a regulated system like DSM.")
add_footer(s3)


# ════════════════════════════════════════════════════════════════
# SLIDE 4 — WITH IBM BOB  (0:52 – 1:12)
# ════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
add_solid_background(s4, SLIDE_BG)
add_rect(s4, 0, 0, 13.33, 0.18, fill_color=GREEN)

add_label(s4, "With IBM Bob \u2014 The Delegate Approach", 0.3, 0.28, 9.0, 0.52,
          font_size=24, bold=True, color=GREEN)
add_label(s4, "ICA Context loaded via MCP \u2192 Bob generated code, tests, docs \u2014 all consistent, all traced",
          0.3, 0.82, 12.0, 0.35, font_size=12, color=MUTED)
divider(s4, 1.2, color=GREEN)

# Flow chain
chain = [
    ("ICA Context\n(MCP)\n17 docs", IBM_BLUE),
    ("Mapping\nRules\nTR-*/VR-*", NG_GREEN),
    ("Java Code\nSpring Batch\n+ REST API", GREEN),
    ("Unit Tests\n44 tests\n0 failures", ACCENT),
    ("React\nDashboard\n+ Docs", RGBColor(0x00, 0x53, 0x9A)),
]
for i, (lbl, col) in enumerate(chain):
    lx = 0.3 + i * 2.6
    add_rect(s4, lx, 1.3, 2.35, 1.45, fill_color=col)
    add_label(s4, lbl, lx + 0.08, 1.38, 2.2, 1.28,
              font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(chain) - 1:
        add_label(s4, "\u2192", lx + 2.37, 1.8, 0.22, 0.5,
                  font_size=20, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

# Metric tiles
metrics = [
    ("~2", "Weeks total\n(was 6-8)", GREEN),
    ("44",  "Unit tests\n0 failures", GREEN),
    ("55",  "Data acceptance\nchecks PASS", NG_GREEN),
    ("17",  "ICA context\ndocs generated", IBM_BLUE),
    ("1",   "Person\n(was 3)", ACCENT),
    ("0",   "Secrets in\n250+ files", RGBColor(0x00, 0x53, 0x9A)),
]
tile_xs = [0.3, 2.47, 4.64, 6.81, 8.98, 11.15]
for i, (val, lbl, col) in enumerate(metrics):
    lx = tile_xs[i]
    add_rect(s4, lx, 2.9, 1.97, 1.45,
             fill_color=GREEN_LIGHT, line_color=col, line_width=Pt(1.5))
    add_label(s4, val, lx + 0.05, 2.98, 1.87, 0.68,
              font_size=34, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_label(s4, lbl, lx + 0.05, 3.6, 1.87, 0.68,
              font_size=9, color=DARK, align=PP_ALIGN.CENTER)

# What Bob generated
add_label(s4, "What Bob Generated for DSM", 0.3, 4.5, 5.0, 0.38,
          font_size=12, bold=True, color=IBM_BLUE)
generated = [
    "Spring Batch pipeline: InitialLoadJob, DailyIncrementalJob, MonthlyUsageJob \u2014 all 8 DSM entity types",
    "Oracle DDL: Flyway V001-V004 migrations for DSM staging/target tables \u2014 idempotent, versioned",
    "Snowflake provisioning: 13 scripts with aborting pre-flight checks and 55-check acceptance suite",
    "REST API: 9 endpoints (job trigger, status, errors, reconciliation, health) + Swagger docs",
    "React + TypeScript dashboard: job control panel, health badges, reconciliation summary, error table",
    "25+ documentation files: architecture, ADRs, security design, operations runbook \u2014 always in sync",
]
for i, g in enumerate(generated):
    add_label(s4, "\u25b8  " + g, 0.3, 4.92 + i * 0.36, 12.7, 0.33,
              font_size=9.5, color=DARK)

script_bar(s4, "0:52-1:12",
           "With Bob at the Delegate stage, I loaded 17 ICA context docs via MCP. "
           "Bob generated the full Spring Batch pipeline, Flyway DDL, Snowflake scripts, REST API, "
           "and React dashboard \u2014 all traced back to those same mapping rules. "
           "6-8 weeks became 2 weeks. 3 people became 1.")
add_footer(s4)


# ════════════════════════════════════════════════════════════════
# SLIDE 5 — QUANTITATIVE COMPARISON  (0:52 – 1:12 continued)
# ════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank)
add_solid_background(s5, SLIDE_BG)
add_rect(s5, 0, 0, 13.33, 0.18, fill_color=IBM_BLUE)

add_label(s5, "Before vs After \u2014 Quantitative", 0.3, 0.28, 9.0, 0.52,
          font_size=24, bold=True, color=IBM_BLUE)
add_label(s5, "Every claim has a number", 0.3, 0.82, 9.0, 0.35, font_size=13, color=MUTED)
divider(s5, 1.2)

# Header row
add_rect(s5, 0.25, 1.28, 4.5, 0.42, fill_color=MUTED)
add_rect(s5, 4.77, 1.28, 3.9, 0.42, fill_color=RED)
add_rect(s5, 8.69, 1.28, 4.4, 0.42, fill_color=GREEN)
for lx, lbl, col in [(0.35, "Activity", WHITE),
                      (4.87, "Without Bob (Manual)", WHITE),
                      (8.79, "With IBM Bob", WHITE)]:
    add_label(s5, lbl, lx, 1.3, 4.3 if lx == 0.35 else (3.7 if lx == 4.87 else 4.3), 0.38,
              font_size=10, bold=True, color=col,
              align=PP_ALIGN.LEFT if lx == 0.35 else PP_ALIGN.CENTER)

rows = [
    ("Requirements & DSM-Snowflake\nmapping analysis",
     "3-4 days  |  BA interviews\nspreadsheet maps, no code link",
     "1 day  |  17 ICA docs, machine-readable\nYAML mapping catalogue"),
    ("Oracle DDL & Flyway migrations",
     "2-3 days  |  hand-written\nenvironment drift risk",
     "2 hours  |  V001-V004, all FK indexes\nidempotent, auto-applied on startup"),
    ("Snowflake provisioning (13 scripts)",
     "3-4 days  |  no pre-flight safety\nerrors need manual cleanup",
     "0.5 days  |  aborting pre-flight checks\n55 acceptance checks, FAIL_COUNT=0"),
    ("Spring Batch pipeline (3 jobs,\n8 DSM entity types)",
     "10-14 days  |  watermark bugs\nretry/skip non-trivial",
     "3-4 days  |  all entities, ICA-traced\nWatermarkService, FkResolution"),
    ("REST API + React Dashboard",
     "8-11 days  |  separate frontend\nresource needed",
     "1.5 days  |  typed Axios client\nRecharts, live polling"),
    ("Unit tests + Documentation",
     "4-5 days  |  deferred, gaps\nonboarding: 2-3 weeks",
     "0.5 days  |  44 tests, 0 failures\nonboarding: 30 minutes"),
    ("TOTAL",
     "6-8 WEEKS  |  3 people",
     "~2 WEEKS  |  1 person  (75% faster)"),
]
for i, (act, before, after) in enumerate(rows):
    ly = 1.72 + i * 0.47
    is_total = (i == len(rows) - 1)
    row_bg   = RGBColor(0x26, 0x26, 0x26) if is_total else (RGBColor(0xFA, 0xFA, 0xFA) if i % 2 == 0 else WHITE)
    bef_bg   = RED   if is_total else RED_LIGHT
    aft_bg   = GREEN if is_total else GREEN_LIGHT
    for lx, w, bg in [(0.25, 4.5, row_bg), (4.77, 3.9, bef_bg), (8.69, 4.4, aft_bg)]:
        add_rect(s5, lx, ly, w, 0.45, fill_color=bg,
                 line_color=RGBColor(0xE0, 0xE0, 0xE0), line_width=Pt(0.5))
    tc    = WHITE if is_total else DARK
    tc_r  = WHITE if is_total else RED
    tc_g  = WHITE if is_total else GREEN
    add_label(s5, act,    0.35, ly + 0.03, 4.3, 0.4, font_size=8.5, bold=is_total, color=tc)
    add_label(s5, before, 4.87, ly + 0.03, 3.7, 0.4, font_size=8.5, bold=is_total, color=tc_r, align=PP_ALIGN.CENTER)
    add_label(s5, after,  8.79, ly + 0.03, 4.2, 0.4, font_size=8.5, bold=is_total, color=tc_g, align=PP_ALIGN.CENTER)

add_footer(s5)


# ════════════════════════════════════════════════════════════════
# SLIDE 6 — ARCHITECTURE DIAGRAM  (1:12 – 1:28)
# ════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank)
add_solid_background(s6, SLIDE_BG)
add_rect(s6, 0, 0, 13.33, 0.18, fill_color=IBM_BLUE)

add_label(s6, "Solution Architecture", 0.3, 0.28, 9.0, 0.52,
          font_size=24, bold=True, color=IBM_BLUE)
add_label(s6,
          "Snowflake (CDP)  \u2192  Spring Batch  \u2192  DSM Oracle Tables  |  React Dashboard  |  ICA Governance",
          0.3, 0.82, 12.5, 0.35, font_size=12, color=MUTED)
divider(s6, 1.2)

# Source zone
add_rect(s6, 0.2, 1.3, 2.6, 3.4,
         fill_color=YELLOW_LIGHT, line_color=RGBColor(0xD9, 0x77, 0x06), line_width=Pt(1))
add_label(s6, "SOURCE", 0.3, 1.35, 2.4, 0.28,
          font_size=8, bold=True, color=RGBColor(0x85, 0x4D, 0x0E), align=PP_ALIGN.CENTER)
add_rect(s6, 0.35, 1.7, 2.3, 0.9,
         fill_color=WHITE, line_color=RGBColor(0xD9, 0x77, 0x06), line_width=Pt(0.75))
add_label(s6, "Snowflake\nCDP_UTIL_DB", 0.35, 1.76, 2.3, 0.55,
          font_size=10, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_label(s6, "CUSTOMER / BILLING\nSTAGING export views", 0.35, 2.32, 2.3, 0.45,
          font_size=8, color=MUTED, align=PP_ALIGN.CENTER)
add_rect(s6, 0.35, 2.85, 2.3, 0.7,
         fill_color=WHITE, line_color=RGBColor(0xD9, 0x77, 0x06), line_width=Pt(0.75))
add_label(s6, "CSS \u2192 Kraken (future)\n\u2192 CDP Snowflake", 0.35, 2.92, 2.3, 0.55,
          font_size=8.5, color=DARK, align=PP_ALIGN.CENTER)

# App Server zone
add_rect(s6, 3.15, 1.3, 6.5, 3.4,
         fill_color=IBM_BLUE_LIGHT, line_color=IBM_BLUE, line_width=Pt(1))
add_label(s6, "APPLICATION SERVER", 3.25, 1.35, 6.3, 0.28,
          font_size=8, bold=True, color=IBM_BLUE, align=PP_ALIGN.CENTER)

add_rect(s6, 3.35, 1.72, 2.5, 0.72, fill_color=IBM_BLUE)
add_label(s6, "Spring Boot REST API\n:8080 | Swagger", 3.35, 1.78, 2.5, 0.6,
          font_size=9, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(s6, 6.05, 1.72, 2.5, 0.72, fill_color=IBM_BLUE)
add_label(s6, "Scheduler\n@Scheduled cron\nInitial/Daily/Monthly", 6.05, 1.76, 2.5, 0.66,
          font_size=8.5, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(s6, 3.35, 2.6, 6.1, 1.9,
         fill_color=WHITE, line_color=IBM_BLUE, line_width=Pt(0.75))
add_label(s6, "Spring Batch Engine", 3.35, 2.63, 6.1, 0.32,
          font_size=10, bold=True, color=IBM_BLUE, align=PP_ALIGN.CENTER)
add_label(s6, "InitialLoadJob  \u00b7  DailyIncrementalJob  \u00b7  MonthlyUsageJob",
          3.35, 2.93, 6.1, 0.28, font_size=8.5, color=MUTED, align=PP_ALIGN.CENTER)
add_label(s6, "ItemReaders (Snowflake JDBC)  \u2192  ItemProcessors (ICA rules)  \u2192  ItemWriters (Oracle MERGE)",
          3.35, 3.2, 6.1, 0.28, font_size=8, color=DARK, align=PP_ALIGN.CENTER)
add_label(s6, "WatermarkService  \u00b7  FkResolutionService  \u00b7  SkipPolicy \u2192 ETL_RECORD_ERROR",
          3.35, 3.48, 6.1, 0.28, font_size=8, color=DARK, align=PP_ALIGN.CENTER)
add_label(s6, "Flyway V001-V004 DDL (auto-applied on startup)",
          3.35, 3.78, 6.1, 0.28, font_size=8, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

# Target zone
add_rect(s6, 10.0, 1.3, 3.1, 3.4,
         fill_color=GREEN_LIGHT, line_color=GREEN, line_width=Pt(1))
add_label(s6, "TARGET (DSM)", 10.1, 1.35, 2.9, 0.28,
          font_size=8, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_rect(s6, 10.15, 1.7, 2.8, 0.72,
         fill_color=WHITE, line_color=GREEN, line_width=Pt(0.75))
add_label(s6, "Oracle 23ai\nDSM Schema", 10.15, 1.76, 2.8, 0.6,
          font_size=10, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_rect(s6, 10.15, 2.52, 2.8, 0.72,
         fill_color=WHITE, line_color=GREEN, line_width=Pt(0.75))
add_label(s6, "TGT_CUSTOMER\nTGT_ENERGY_ACCOUNT\n+5 DSM tables", 10.15, 2.56, 2.8, 0.65,
          font_size=8.5, color=DARK, align=PP_ALIGN.CENTER)
add_rect(s6, 10.15, 3.35, 2.8, 0.62,
         fill_color=WHITE, line_color=GREEN, line_width=Pt(0.75))
add_label(s6, "ETL_WATERMARK\nETL_JOB_RUN\nETL_RECONCILIATION", 10.15, 3.38, 2.8, 0.58,
          font_size=8, color=DARK, align=PP_ALIGN.CENTER)

# Frontend zone
add_rect(s6, 0.2, 4.85, 2.6, 1.2,
         fill_color=RGBColor(0xED, 0xE9, 0xFE), line_color=ACCENT, line_width=Pt(1))
add_label(s6, "FRONTEND", 0.3, 4.9, 2.4, 0.28,
          font_size=8, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_label(s6, "React + TypeScript\nVite :5173\nRecharts Dashboard", 0.3, 5.2, 2.4, 0.75,
          font_size=9, color=DARK, align=PP_ALIGN.CENTER)

# Governance bar
add_rect(s6, 0.2, 6.18, 12.9, 0.52,
         fill_color=RGBColor(0xF4, 0xF4, 0xF4),
         line_color=RGBColor(0xC6, 0xC6, 0xC6), line_width=Pt(0.75))
add_label(s6,
          "GOVERNANCE (Design-time)  \u2014  IBM Bob / ICA Context Studio (MCP)  "
          "\u00b7  17 docs: GL, CM, TR, VR, Error, Recon rules  "
          "\u2192  code  \u2192  tests  \u2192  docs",
          0.4, 6.24, 12.6, 0.42,
          font_size=9, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

# Arrow labels
add_label(s6, "JDBC read \u2192", 2.82, 3.02, 0.38, 0.28,
          font_size=8, bold=True, color=RGBColor(0xD9, 0x77, 0x06))
add_label(s6, "\u2192 MERGE", 9.45, 3.02, 0.6, 0.28,
          font_size=8, bold=True, color=GREEN)
add_label(s6, "REST", 1.62, 4.6, 1.5, 0.28,
          font_size=8, color=ACCENT, align=PP_ALIGN.CENTER)

script_bar(s6, "1:12-1:28",
           "Here is the architecture: Snowflake as source, Spring Batch as the engine, "
           "Oracle DSM tables as the target. The ICA context studio feeds Bob with mapping rules at design time. "
           "The running application has zero ICA dependency \u2014 it is all compiled in.")
add_footer(s6)


# ════════════════════════════════════════════════════════════════
# SLIDE 7 — TRACEABILITY & KEY DIFFERENTIATOR  (1:28 – 1:45)
# ════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank)
add_solid_background(s7, SLIDE_BG)
add_rect(s7, 0, 0, 13.33, 0.18, fill_color=ACCENT)

add_label(s7, "The Key Differentiator \u2014 Traceability", 0.3, 0.28, 9.0, 0.52,
          font_size=24, bold=True, color=ACCENT)
add_label(s7, "ICA Rule \u2192 Java Class \u2192 Unit Test  \u00b7  Zero drift between spec and code on a regulated system",
          0.3, 0.82, 12.0, 0.35, font_size=12, color=MUTED)
divider(s7, 1.2, color=ACCENT)

# Chain
chain2 = [
    ("ICA Context\nStudio (MCP)\n17 docs", IBM_BLUE),
    ("Transformation\nRules TR-*\nValidation VR-*", NG_GREEN),
    ("Java Classes\nProcessors &\nValidators", GREEN),
    ("44 Unit Tests\n0 Failures\nAll ICA-traced", ACCENT),
    ("DSM Oracle\nMERGE +\nAudit trail", RGBColor(0x00, 0x53, 0x9A)),
]
for i, (lbl, col) in enumerate(chain2):
    lx = 0.3 + i * 2.6
    add_rect(s7, lx, 1.32, 2.35, 1.4, fill_color=col)
    add_label(s7, lbl, lx + 0.08, 1.4, 2.2, 1.22,
              font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(chain2) - 1:
        add_label(s7, "\u2192", lx + 2.38, 1.82, 0.22, 0.5,
                  font_size=20, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

# Example trace box
add_rect(s7, 0.3, 2.9, 12.7, 1.62,
         fill_color=WHITE, line_color=ACCENT, line_width=Pt(1))
add_label(s7, "Example: TR-CUST-STATUS-01  (DSM Customer Status Mapping)", 0.5, 2.97, 9.0, 0.35,
          font_size=12, bold=True, color=ACCENT)
trace = [
    ("ICA Rule",   "TR-CUST-STATUS-01 in 06-transformation-rules.md: "
                   "Snowflake ACCOUNT_STATUS='ACTIVE' \u2192 Oracle IS_ACTIVE=1, inactive \u2192 0"),
    ("Java Class", "CustomerProcessor.java \u2014 mapAccountStatus() method, generated directly from ICA rule ID"),
    ("Unit Test",  "CustomerProcessorTest.java \u2014 testActiveStatusMapsToOne(), testInactiveStatusMapsToZero()"),
    ("Matrix",     "docs/ica-mcp-validation.md \u2014 full rule \u2194 class \u2194 test traceability matrix for 30+ DSM rules"),
]
for i, (lbl, val) in enumerate(trace):
    ly = 3.35 + i * 0.27
    add_label(s7, lbl + ":", 0.5, ly, 1.4, 0.25,
              font_size=9.5, bold=True, color=MUTED)
    add_label(s7, val, 2.0, ly, 10.8, 0.25,
              font_size=9.5, color=DARK)

# Why it matters for DSM
add_label(s7, "Why Traceability Matters on a Regulated DSM System", 0.3, 4.68, 8.0, 0.35,
          font_size=12, bold=True, color=IBM_BLUE)
points = [
    "DSM reports to MA & NY regulators \u2014 mapping errors in customer data affect programme compliance and audit outcomes.",
    "Business rules live in ICA, not a developer's head. When Kraken data schema changes, Bob regenerates only affected code.",
    "Every VR-*/TR-* rule has a test. Regression caught automatically \u2014 not discovered in a regulatory submission.",
    "Onboarding a new DSM engineer: 30 minutes (README + ICA bundle). Previously 2-3 weeks of tribal knowledge transfer.",
]
for i, pt in enumerate(points):
    add_label(s7, "\u25b8  " + pt, 0.3, 5.08 + i * 0.36, 12.7, 0.33,
              font_size=10, color=DARK)

script_bar(s7, "1:28-1:45",
           "The real differentiator is not speed \u2014 it is consistency. "
           "DSM feeds regulated energy efficiency programmes. A mapping error is not just a bug, it is a compliance risk. "
           "Bob maintained spec-to-code traceability across 80 backlog items and 6 phases. "
           "That discipline normally needs a dedicated BA and QA team.")
add_footer(s7)


# ════════════════════════════════════════════════════════════════
# SLIDE 8 — CALL TO ACTION  (1:45 – 2:00)
# ════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(blank)
add_solid_background(s8, SLIDE_BG_DARK)
add_rect(s8, 0, 0, 13.33, 0.18, fill_color=NG_GREEN)

add_label(s8, "What This Means for National Grid", 0.4, 0.28, 10.0, 0.6,
          font_size=26, bold=True, color=WHITE)
add_label(s8, "And the next step to get there",
          0.4, 0.9, 8.0, 0.35, font_size=14, color=NG_GREEN_LIGHT)
divider(s8, 1.3, color=NG_GREEN)

# Value banner
add_rect(s8, 0.3, 1.4, 12.7, 1.0, fill_color=IBM_BLUE)
add_label(s8,
          "IBM Bob compressed a 6-8 week, 3-person integration project into ~2 weeks with 1 person \u2014\n"
          "with full traceability, 44 tests passing, and a production-ready Snowflake \u2192 DSM loader prototype.",
          0.5, 1.47, 12.3, 0.9,
          font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# CTA cards
ctas = [
    ("30-min\nWorkshop",
     "Live demo: Bob generates a\nDSM mapping doc + processor\nfrom a new ICA rule.\nQ&A on AI Maturity Curve.",
     IBM_BLUE),
    ("Proof of\nExperience (PoX)",
     "Extend the DSM prototype:\nadd the Kraken-to-Snowflake\npipeline or another DSM\nsubprogramme in 1-2 weeks.",
     NG_GREEN),
    ("Deeper Dive\nMeeting",
     "ICA Context Studio + MCP\nwalkthrough for National Grid's\narchitect / data lead.\nRoadmap to Scale & Extend.",
     ACCENT),
]
for i, (title, desc, col) in enumerate(ctas):
    lx = 0.3 + i * 4.36
    add_rect(s8, lx, 2.58, 4.1, 3.0, fill_color=col)
    add_label(s8, title, lx + 0.1, 2.72, 3.9, 0.9,
              font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s8, lx + 0.15, 3.62, 3.8, 0.04, fill_color=WHITE)
    add_label(s8, desc, lx + 0.15, 3.72, 3.8, 1.7,
              font_size=11, color=WHITE, align=PP_ALIGN.CENTER)

# Key metrics recap
add_multiline(s8, [
    ("KEY METRICS RECAP", True, IBM_BLUE_LIGHT),
    ("", False, WHITE),
    ("6-8 weeks  \u2192  ~2 weeks    \u00b7    3 people  \u2192  1 person    \u00b7    44 tests, 0 failures    \u00b7    "
     "55 acceptance checks, FAIL_COUNT = 0", False, WHITE),
    ("17 ICA context docs    \u00b7    0 secrets in 250+ files    \u00b7    "
     "30 min developer onboarding    \u00b7    Full regulatory traceability", False, WHITE),
], 0.3, 5.78, 12.7, 1.0, font_size=12, align=PP_ALIGN.CENTER)

# Script note
add_rect(s8, 0.3, 6.9, 12.7, 0.35,
         fill_color=RGBColor(0x39, 0x39, 0x3A),
         line_color=RGBColor(0x52, 0x52, 0x52), line_width=Pt(0.5))
add_label(s8,
          "SCRIPT [1:45-2:00]  "
          "\"For National Grid's DSM and the Kraken migration ahead, IBM Bob means your engineers focus "
          "on strategy and integration design, while Bob handles implementation. "
          "Let's set up a 30-min workshop. Thank you.\"",
          0.45, 6.92, 12.5, 0.32,
          font_size=8, italic=True, color=RGBColor(0xC6, 0xC6, 0xC6))

add_footer(s8, "National Grid DSM  |  CDP Snowflake Loader  |  IBM Bob Level 3  |  Althaf P M  |  2025  |  Made with IBM Bob")


# ── Save ─────────────────────────────────────────────────────────
out_dir  = os.path.join(os.path.dirname(__file__), '..', 'docs')
os.makedirs(out_dir, exist_ok=True)
out_path = os.path.join(out_dir, 'IBM_Bob_Stand_Deliver.pptx')
prs.save(out_path)
print("SAVED: " + os.path.abspath(out_path))
